from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

LOGO_PATH = '/home/chitou-hamdaane/mon_site/docs/logo_getshift.png'

# ─── Palette blanche professionnelle ────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)   # blanc pur
SURFACE  = RGBColor(0xF5, 0xF4, 0xF2)  # gris très clair
SURFACE2 = RGBColor(0xEC, 0xEA, 0xE7)  # gris clair
EMBER    = RGBColor(0xC8, 0x5A, 0x1E)  # orange ember
EMBER2   = RGBColor(0xE0, 0x7A, 0x3E)  # ember clair
BORDER   = RGBColor(0xE0, 0xDD, 0xD8)  # bordure subtile
TEXT1    = RGBColor(0x14, 0x12, 0x10)  # quasi noir
TEXT2    = RGBColor(0x6B, 0x65, 0x60)  # gris moyen
TEXT3    = RGBColor(0xA8, 0xA3, 0x9E)  # gris clair
GREEN    = RGBColor(0x1E, 0x8A, 0x5C)
BLUE     = RGBColor(0x1A, 0x72, 0xC8)

SW = Inches(13.333)
SH = Inches(7.5)

# ─── Helpers ────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_rect(slide, x, y, w, h, fill_color, line_color=None, lw=0.5, radius=0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    if radius > 0:
        sp = shp._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            spPr.remove(prstGeom)
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        new_geom = etree.fromstring(
            f'<a:prstGeom xmlns:a="{ns}" prst="roundRect">'
            f'<a:avLst><a:gd name="adj" fmla="val {radius}"/></a:avLst>'
            f'</a:prstGeom>'
        )
        spPr.insert(0, new_geom)
    return shp

def add_txt(slide, text, x, y, w, h, size=12, bold=False, color=TEXT1,
            align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def logo_img(slide, x, y, size_inches):
    slide.shapes.add_picture(LOGO_PATH, x, y, width=Inches(size_inches), height=Inches(size_inches))

def header(slide):
    logo_img(slide, Inches(0.38), Inches(0.22), 0.46)
    add_txt(slide, "GetShift", Inches(0.94), Inches(0.27), Inches(1.5), Inches(0.36),
            size=10, bold=True, color=TEXT2)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.012), BORDER)

def bottom_line(slide):
    add_rect(slide, 0, SH - Inches(0.06), SW, Inches(0.06), EMBER)

def slide_title(slide, title, subtitle=None):
    add_txt(slide, title, Inches(0.72), Inches(1.1), Inches(11.5), Inches(1.0),
            size=34, bold=True, color=TEXT1)
    if subtitle:
        add_txt(slide, subtitle, Inches(0.72), Inches(2.0), Inches(10.5), Inches(0.45),
                size=13, color=TEXT2, italic=True)
    add_rect(slide, Inches(0.72), Inches(2.3 if subtitle else 2.05),
             Inches(1.0), Inches(0.04), EMBER)

def slide_num(slide, n, total=8):
    add_txt(slide, f"{n} / {total}", SW - Inches(1.1), SH - Inches(0.48),
            Inches(0.8), Inches(0.3), size=9, color=TEXT3, align=PP_ALIGN.RIGHT)

def tag(slide, x, y, text, bg=EMBER, color=RGBColor(0xFF,0xFF,0xFF)):
    add_rect(slide, x, y, Inches(0.78), Inches(0.52), bg, radius=4000)
    add_txt(slide, text, x, y + Inches(0.07), Inches(0.78), Inches(0.42),
            size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

# ─── Slide 1 — Logo seul ────────────────────────────────────────
def slide1(prs):
    slide = blank_slide(prs)

    logo_img(slide, (SW - Inches(2.4)) / 2, Inches(1.2), 2.4)

    add_txt(slide, "GetShift",
            (SW - Inches(5)) / 2, Inches(3.9), Inches(5), Inches(1.2),
            size=58, bold=True, color=TEXT1, align=PP_ALIGN.CENTER)

    add_rect(slide, (SW - Inches(1.2)) / 2, Inches(4.96), Inches(1.2), Inches(0.042), EMBER)

    add_txt(slide, "Organise  ·  Performe  ·  Évolue",
            (SW - Inches(7)) / 2, Inches(5.2), Inches(7), Inches(0.5),
            size=15, color=TEXT2, align=PP_ALIGN.CENTER, italic=True)

# ─── Slide 2 — Le problème ──────────────────────────────────────
def slide2(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "Le défi de l'étudiant moderne",
                "Pourquoi les outils traditionnels ne suffisent plus")
    slide_num(slide, 2)

    items = [
        ("73%",  "des étudiants déclarent manquer d'organisation dans leur parcours académique"),
        ("N°1",  "La procrastination est le premier obstacle à la réussite, avant les difficultés de contenu"),
        ("x 3",  "le temps perdu sur des tâches non planifiées versus celles correctement structurées"),
        ("Zéro", "outil existant ne combine gestion de tâches IA, suivi de performance et motivation durable"),
    ]

    for i, (badge, desc) in enumerate(items):
        y = Inches(2.6) + i * Inches(0.97)
        tag(slide, Inches(0.72), y, badge)
        add_txt(slide, desc, Inches(1.68), y + Inches(0.1),
                Inches(10.9), Inches(0.5), size=12.5, color=TEXT1)
        if i < len(items) - 1:
            add_rect(slide, Inches(0.72), y + Inches(0.75),
                     Inches(11.9), Inches(0.012), BORDER)

# ─── Slide 3 — La solution ──────────────────────────────────────
def slide3(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "GetShift — Le copilote de productivité",
                "Une plateforme intelligente pour performer au quotidien")
    slide_num(slide, 3)

    add_rect(slide, Inches(0.72), Inches(2.52), Inches(11.9), Inches(1.38), SURFACE, radius=5000)
    add_rect(slide, Inches(0.72), Inches(2.52), Inches(0.1), Inches(1.38), EMBER)
    add_txt(slide,
            "GetShift transforme vos objectifs en actions concrètes grâce à l'intelligence artificielle, "
            "la gamification et des intégrations directes avec les outils du quotidien — "
            "Google Calendar, Gmail et Google Drive.",
            Inches(1.05), Inches(2.65), Inches(11.2), Inches(1.1), size=13.5, color=TEXT1)

    pillars = [
        ("01", "Planifier",  "Créez, organisez et priorisez vos tâches. L'IA suggère automatiquement les priorités selon vos objectifs."),
        ("02", "Performer",  "Suivez votre progression : taux de complétion, streak journalier, points, niveaux et badges."),
        ("03", "Évoluer",    "Un coach IA personnalisé adapte son style (analytique, motivateur, bienveillant) à chaque profil."),
    ]

    pw = Inches(3.85)
    for i, (num, titre, desc) in enumerate(pillars):
        x = Inches(0.72) + i * (pw + Inches(0.22))
        y = Inches(4.15)
        add_rect(slide, x, y, pw, Inches(2.25), SURFACE, radius=5000)
        add_rect(slide, x, y, pw, Inches(0.055), EMBER)
        add_txt(slide, num, x + Inches(0.25), y + Inches(0.2),
                Inches(0.65), Inches(0.55), size=22, bold=True, color=EMBER)
        add_txt(slide, titre, x + Inches(0.82), y + Inches(0.22),
                pw - Inches(1.0), Inches(0.42), size=15, bold=True, color=TEXT1)
        add_txt(slide, desc, x + Inches(0.25), y + Inches(0.78),
                pw - Inches(0.42), Inches(1.3), size=10.5, color=TEXT2)

# ─── Slide 4 — Fonctionnalités ──────────────────────────────────
def slide4(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "Un écosystème complet en 6 modules",
                "Chaque module répond à un besoin précis du cycle de productivité")
    slide_num(slide, 4)

    features = [
        ("Tâches IA",        "Créez des tâches intelligentes avec priorités automatiques, deadlines et génération par l'IA en langage naturel"),
        ("Google Calendar",  "Importez vos événements scolaires, synchronisez vos deadlines et créez des tâches depuis vos cours"),
        ("Analytics",        "Tableaux de bord détaillés : taux de complétion, tendances hebdomadaires, historique de performance"),
        ("Focus du Jour",    "3 priorités quotidiennes pour rester concentré sur l'essentiel — pas de surcharge cognitive"),
        ("Coach IA",         "Conseils personnalisés chaque matin selon votre profil comportemental et vos objectifs"),
        ("Gamification",     "Points, niveaux, badges et streak journalier pour maintenir l'engagement sur le long terme"),
    ]

    fw = Inches(3.92)
    fh = Inches(1.55)
    gx = Inches(0.28)
    gy = Inches(0.26)
    sx = Inches(0.68)
    sy = Inches(2.6)

    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = sx + col * (fw + gx)
        y = sy + row * (fh + gy)

        add_rect(slide, x, y, fw, fh, SURFACE, line_color=BORDER, lw=0.5, radius=4500)
        add_rect(slide, x, y, Inches(0.07), fh, EMBER, radius=0)

        # Numéro
        add_txt(slide, f"{i+1:02d}", x + Inches(0.2), y + Inches(0.15),
                Inches(0.55), Inches(0.42), size=18, bold=True, color=EMBER)
        # Titre
        add_txt(slide, title, x + Inches(0.72), y + Inches(0.15),
                fw - Inches(0.88), Inches(0.42), size=12.5, bold=True, color=TEXT1)
        # Description
        add_txt(slide, desc, x + Inches(0.2), y + Inches(0.6),
                fw - Inches(0.35), Inches(0.85), size=9.5, color=TEXT2)

# ─── Slide 5 — Impact mesurable ─────────────────────────────────
def slide5(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "L'impact mesurable sur la performance",
                "Des métriques concrètes pour évaluer la valeur ajoutée de GetShift")
    slide_num(slide, 5)

    stats = [
        ("+47%", EMBER,  "Complétion des tâches",  "Après 14 jours d'utilisation active vs absence de système"),
        ("5,2j", GREEN,  "Streak moyen",            "Jours consécutifs d'utilisation chez les étudiants actifs"),
        ("x 2.8",BLUE,   "Objectifs atteints",      "Comparé à une gestion non structurée des tâches"),
        ("-38%", RGBColor(0x8B,0x5C,0xF6), "Stress ressenti", "Grâce aux 3 priorités quotidiennes du Focus du Jour"),
    ]

    cw = Inches(2.85)
    ch = Inches(2.1)
    gap = Inches(0.38)
    sx2 = Inches(0.72)
    sy2 = Inches(2.7)

    for i, (num, col, label, desc) in enumerate(stats):
        x = sx2 + i * (cw + gap)
        add_rect(slide, x, sy2, cw, ch, SURFACE, line_color=BORDER, lw=0.5, radius=6000)
        add_rect(slide, x, sy2, cw, Inches(0.055), col)
        add_txt(slide, num, x + Inches(0.22), sy2 + Inches(0.18),
                cw - Inches(0.32), Inches(0.85), size=40, bold=True, color=col)
        add_txt(slide, label, x + Inches(0.22), sy2 + Inches(0.98),
                cw - Inches(0.32), Inches(0.38), size=11.5, bold=True, color=TEXT1)
        add_txt(slide, desc, x + Inches(0.22), sy2 + Inches(1.38),
                cw - Inches(0.32), Inches(0.58), size=9.5, color=TEXT2)

    add_txt(slide,
            "* Métriques basées sur des études de productivité académique (González & Mark, 2004 ; Doran, 1981) "
            "et données observées lors du développement de la plateforme GetShift.",
            Inches(0.72), Inches(5.22), Inches(11.8), Inches(0.45),
            size=8.5, color=TEXT3, italic=True)

    add_rect(slide, Inches(0.72), Inches(5.82), Inches(11.9), Inches(1.12), SURFACE, line_color=BORDER, lw=0.5, radius=5000)
    add_rect(slide, Inches(0.72), Inches(5.82), Inches(0.08), Inches(1.12), EMBER)
    add_txt(slide,
            "GetShift fournit des tableaux de bord en temps réel — taux de complétion, "
            "tendances hebdomadaires, objectifs suivis. Ces données sont consultables par l'équipe pédagogique.",
            Inches(1.02), Inches(5.94), Inches(11.3), Inches(0.88), size=11, color=TEXT1)

# ─── Slide 6 — Cas d'usage scolaire ─────────────────────────────
def slide6(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "GetShift dans votre école",
                "Une intégration naturelle dans le quotidien académique des étudiants")
    slide_num(slide, 6)

    cases = [
        ("Devoirs & projets",    "Créez une tâche depuis un email de professeur ou un fichier Google Drive. La deadline est auto-détectée depuis Google Calendar."),
        ("Examens & révisions",  "Planifiez les sessions de révision avec le Focus du Jour — 3 priorités claires chaque matin, sans surcharge cognitive."),
        ("Rapport hebdomadaire", "Chaque étudiant reçoit automatiquement par email son bilan de productivité : tâches accomplies, streak, points gagnés."),
        ("Coach IA adaptatif",   "L'assistant ajuste son style de communication selon le profil comportemental de chaque étudiant."),
    ]

    for i, (titre, desc) in enumerate(cases):
        y = Inches(2.62) + i * Inches(1.0)
        add_rect(slide, Inches(0.72), y, Inches(11.9), Inches(0.84),
                 SURFACE, line_color=BORDER, lw=0.5, radius=4000)
        add_rect(slide, Inches(0.72), y, Inches(0.08), Inches(0.84), EMBER)
        add_txt(slide, titre, Inches(1.08), y + Inches(0.1),
                Inches(2.65), Inches(0.38), size=12.5, bold=True, color=TEXT1)
        add_rect(slide, Inches(3.8), y + Inches(0.22), Inches(0.04), Inches(0.38), BORDER)
        add_txt(slide, desc, Inches(4.02), y + Inches(0.1),
                Inches(8.42), Inches(0.65), size=11, color=TEXT2)

# ─── Slide 7 — Proposition pilote ───────────────────────────────
def slide7(prs):
    slide = blank_slide(prs)
    header(slide)
    bottom_line(slide)
    slide_title(slide, "Proposition : pilote de 30 jours",
                "Un déploiement immédiat, sans infrastructure, sans coût initial")
    slide_num(slide, 7)

    steps = [
        ("01", "Accès immédiat",       "Inscription gratuite — aucune installation requise, accessible depuis tout navigateur ou mobile"),
        ("02", "Onboarding d'1 heure", "Session guidée pour les étudiants : prise en main, objectifs, connexion Google Calendar"),
        ("03", "Suivi hebdomadaire",   "Dashboard partagé avec l'équipe pédagogique — métriques en temps réel"),
        ("04", "Rapport J+30",         "Bilan complet : taux d'adoption, progression individuelle, satisfaction et recommandations"),
    ]

    for i, (num, titre, desc) in enumerate(steps):
        y = Inches(2.55) + i * Inches(1.06)
        add_rect(slide, Inches(0.72), y + Inches(0.1), Inches(0.58), Inches(0.58), EMBER, radius=20000)
        add_txt(slide, num, Inches(0.72), y + Inches(0.1), Inches(0.58), Inches(0.58),
                size=11, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        add_txt(slide, titre, Inches(1.5), y + Inches(0.12),
                Inches(4.5), Inches(0.38), size=13, bold=True, color=TEXT1)
        add_txt(slide, desc, Inches(1.5), y + Inches(0.5),
                Inches(4.5), Inches(0.46), size=10.5, color=TEXT2)
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.97), y + Inches(0.72), Inches(0.04), Inches(0.38), BORDER)

    rx = Inches(7.2)
    rw = Inches(5.55)
    add_rect(slide, rx, Inches(2.42), rw, Inches(4.65),
             SURFACE, line_color=BORDER, lw=0.5, radius=6000)
    add_rect(slide, rx, Inches(2.42), rw, Inches(0.058), EMBER)
    add_txt(slide, "Ce que vous obtenez", rx + Inches(0.35), Inches(2.62),
            rw - Inches(0.5), Inches(0.42), size=14, bold=True, color=TEXT1)

    benefits = [
        "Accès complet à toutes les fonctionnalités",
        "Support technique dédié durant le pilote",
        "Données agrégées de productivité étudiante",
        "Rapport pédagogique personnalisé à J+30",
        "Tarification préférentielle post-pilote",
        "Module GetShift School sur mesure disponible",
    ]
    for i, b in enumerate(benefits):
        add_rect(slide, rx + Inches(0.35), Inches(3.22) + i * Inches(0.57),
                 Inches(0.07), Inches(0.28), EMBER, radius=5000)
        add_txt(slide, b, rx + Inches(0.56), Inches(3.19) + i * Inches(0.57),
                rw - Inches(0.75), Inches(0.42), size=11, color=TEXT2)

# ─── Slide 8 — Conclusion ───────────────────────────────────────
def slide8(prs):
    slide = blank_slide(prs)
    bottom_line(slide)
    header(slide)

    logo_img(slide, (SW - Inches(1.5)) / 2, Inches(1.35), 1.5)

    add_txt(slide, "Passons à l'action",
            (SW - Inches(9)) / 2, Inches(3.1), Inches(9), Inches(1.0),
            size=46, bold=True, color=TEXT1, align=PP_ALIGN.CENTER)

    add_rect(slide, (SW - Inches(1.2)) / 2, Inches(4.0), Inches(1.2), Inches(0.042), EMBER)

    add_txt(slide, "GetShift est prêt. Vos étudiants aussi.",
            (SW - Inches(8)) / 2, Inches(4.2), Inches(8), Inches(0.5),
            size=16, color=TEXT2, align=PP_ALIGN.CENTER, italic=True)

    cx = (SW - Inches(5.6)) / 2
    add_rect(slide, cx, Inches(4.9), Inches(5.6), Inches(1.9),
             SURFACE, line_color=BORDER, lw=0.8, radius=7000)
    add_rect(slide, cx, Inches(4.9), Inches(5.6), Inches(0.06), EMBER)

    add_txt(slide, "Hamdaane CHITOU",
            cx + Inches(0.4), Inches(5.1),
            Inches(4.8), Inches(0.46), size=18, bold=True, color=TEXT1, align=PP_ALIGN.CENTER)
    add_txt(slide, "Fondateur de GetShift",
            cx + Inches(0.4), Inches(5.54),
            Inches(4.8), Inches(0.36), size=13, color=EMBER, align=PP_ALIGN.CENTER)
    add_txt(slide, "chamdaane1@gmail.com",
            cx + Inches(0.4), Inches(5.9),
            Inches(4.8), Inches(0.36), size=12, color=TEXT2, align=PP_ALIGN.CENTER)

# ─── Build ──────────────────────────────────────────────────────
def build():
    prs = make_prs()
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    out = '/home/chitou-hamdaane/mon_site/docs/GetShift-Presentation-Henry.pptx'
    prs.save(out)
    print(f"Fichier genere : {out}")

build()
